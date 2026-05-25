import os
from pathlib import Path
import pandas as pd


SUPPORTED_OUTPUT_FORMATS = {
    '.txt', '.docx', '.xlsx', '.pptx', '.pdf'
}


def validate_output_filename(filename, default_extension='.txt'):
    """Validate output filenames and ensure a supported extension."""
    if not filename or not isinstance(filename, str):
        raise ValueError('Output filename must be a non-empty string.')

    filename = filename.strip()
    if filename == '':
        raise ValueError('Output filename cannot be empty.')

    path = Path(filename)
    if path.name == '' or path.name.startswith('.'):
        raise ValueError('Output filename must contain a valid name before the extension.')

    extension = path.suffix.lower()
    if extension == '':
        path = path.with_suffix(default_extension)
        extension = path.suffix.lower()

    if extension not in SUPPORTED_OUTPUT_FORMATS:
        raise ValueError(f'Unsupported export format {extension}. Supported formats: {sorted(SUPPORTED_OUTPUT_FORMATS)}')

    return path


class ExportManager:
    """Export structured analysis content into files."""

    @staticmethod
    def export_text_report(report_sections, filename):
        path = validate_output_filename(filename, default_extension='.txt')
        with open(path, 'w', encoding='utf-8') as fd:
            for title, content in report_sections.items():
                fd.write(f'{title}\n')
                fd.write(f'{"=" * len(title)}\n')
                fd.write(f'{content}\n\n')
        return path

    @staticmethod
    def export_docx_report(report_sections, filename):
        path = validate_output_filename(filename, default_extension='.docx')
        try:
            from docx import Document
        except ImportError as exc:
            raise ImportError('python-docx is required for DOCX exports') from exc

        document = Document()
        for title, content in report_sections.items():
            document.add_heading(title, level=1)
            for paragraph in str(content).split('\n'):
                document.add_paragraph(paragraph)
            document.add_paragraph('')
        document.save(path)
        return path

    @staticmethod
    def export_xlsx_report(report_sections, filename):
        path = validate_output_filename(filename, default_extension='.xlsx')
        with pd.ExcelWriter(path, engine='openpyxl') as writer:
            for title, content in report_sections.items():
                sheet_name = title[:31]
                lines = str(content).split('\n')
                df = pd.DataFrame({'Text': lines})
                df.to_excel(writer, sheet_name=sheet_name, index=False)
        return path

    @staticmethod
    def export_pptx_report(report_sections, filename):
        path = validate_output_filename(filename, default_extension='.pptx')
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise ImportError('python-pptx is required for PPTX exports') from exc

        presentation = Presentation()
        for title, content in report_sections.items():
            slide_layout = presentation.slide_layouts[1] if len(presentation.slides) else presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
            body_shape = slide.shapes.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()
            for paragraph in str(content).split('\n'):
                p = text_frame.add_paragraph()
                p.text = paragraph
                p.font.size = Pt(12)
        presentation.save(path)
        return path

    @staticmethod
    def export_pdf_report(report_sections, filename):
        path = validate_output_filename(filename, default_extension='.pdf')
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
        except ImportError as exc:
            raise ImportError('reportlab is required for PDF exports') from exc

        styles = getSampleStyleSheet()
        document = SimpleDocTemplate(str(path), pagesize=letter)
        story = []
        for title, content in report_sections.items():
            story.append(Paragraph(title, styles['Heading1']))
            for paragraph in str(content).split('\n'):
                story.append(Paragraph(paragraph, styles['BodyText']))
            story.append(Spacer(1, 12))
        document.build(story)
        return path


def get_mime_type_for_path(path: str) -> str:
    """Return a sensible MIME type for a filename based on its extension."""
    mapping = {
        '.txt': 'text/plain',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.pdf': 'application/pdf',
    }
    ext = Path(path).suffix.lower()
    return mapping.get(ext, 'application/octet-stream')
